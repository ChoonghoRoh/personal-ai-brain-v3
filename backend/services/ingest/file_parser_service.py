"""파일 파서 서비스

Phase 9-4-1: HWP/HWPX 파일 지원 추가
"""
import logging
from typing import Optional, List, Dict
from pathlib import Path
import io

logger = logging.getLogger(__name__)


class FileParserService:
    """파일 파서 서비스 클래스"""

    def __init__(self):
        self.supported_formats = {
            '.md': self.parse_markdown,
            '.txt': self.parse_text,
            '.pdf': self.parse_pdf,
            '.docx': self.parse_docx,
            '.xlsx': self.parse_excel,
            '.xls': self.parse_excel,
            '.pptx': self.parse_powerpoint,
            '.ppt': self.parse_powerpoint,
            '.hwp': self.parse_hwp,
            '.hwpx': self.parse_hwp,  # Phase 9-4-1: HWPX 지원 추가
            '.jpg': self.parse_image_ocr,
            '.jpeg': self.parse_image_ocr,
            '.png': self.parse_image_ocr,
            '.gif': self.parse_image_ocr
        }
    
    def parse_markdown(self, file_path: Path) -> Optional[str]:
        """Markdown 파일 파싱"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Markdown 파싱 실패: {e}")
            return None
    
    def parse_pdf(self, file_path: Path) -> Optional[str]:
        """PDF 파일 파싱"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return "\n\n".join(text_parts)
        except ImportError:
            logger.warning("pypdf가 설치되지 않았습니다.")
            return None
        except Exception as e:
            logger.error(f"PDF 파싱 실패: {e}")
            return None
    
    def parse_docx(self, file_path: Path) -> Optional[str]:
        """DOCX 파일 파싱"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            logger.warning("python-docx가 설치되지 않았습니다.")
            return None
        except Exception as e:
            logger.error(f"DOCX 파싱 실패: {e}")
            return None
    
    def parse_excel(self, file_path: Path) -> Optional[str]:
        """Excel 파일 파싱"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return "\n".join(text_parts)
        except ImportError:
            logger.warning("openpyxl이 설치되지 않았습니다.")
            return None
        except Exception as e:
            logger.error(f"Excel 파싱 실패: {e}")
            return None
    
    def parse_powerpoint(self, file_path: Path) -> Optional[str]:
        """PowerPoint 파일 파싱"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== 슬라이드 {i} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            logger.warning("python-pptx가 설치되지 않았습니다.")
            return None
        except Exception as e:
            logger.error(f"PowerPoint 파싱 실패: {e}")
            return None
    
    def parse_hwp(self, file_path: Path) -> Optional[str]:
        """HWP/HWPX 파일 파싱 (Phase 9-4-1)"""
        try:
            from backend.services.ingest.hwp_parser import get_hwp_parser

            parser = get_hwp_parser()
            result = parser.parse(file_path)

            if result.success:
                logger.info(f"HWP 파싱 성공: {file_path.name} (type: {result.file_type})")
                return result.text
            else:
                logger.warning(f"HWP 파싱 실패: {result.error}")
                return None

        except ImportError as e:
            logger.warning(f"HWP 파서 모듈 로드 실패: {e}")
            return None
        except Exception as e:
            logger.error(f"HWP 파싱 실패: {e}")
            return None

    def parse_text(self, file_path: Path) -> Optional[str]:
        """텍스트 파일 파싱"""
        try:
            # 여러 인코딩 시도
            encodings = ['utf-8', 'cp949', 'euc-kr', 'utf-16']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            logger.warning(f"텍스트 파일 인코딩 감지 실패: {file_path}")
            return None
        except Exception as e:
            logger.error(f"텍스트 파일 파싱 실패: {e}")
            return None
    
    def parse_image_ocr(self, file_path: Path) -> Optional[str]:
        """이미지 OCR 파싱"""
        try:
            from PIL import Image
            import pytesseract
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='kor+eng')
            return text
        except ImportError:
            logger.warning("PIL 또는 pytesseract가 설치되지 않았습니다.")
            return None
        except Exception as e:
            logger.error(f"이미지 OCR 실패: {e}")
            return None
    
    def parse_file(self, file_path: Path) -> Optional[str]:
        """파일 파싱 (자동 형식 감지)"""
        suffix = file_path.suffix.lower()
        
        if suffix in self.supported_formats:
            parser = self.supported_formats[suffix]
            return parser(file_path)
        else:
            logger.warning(f"지원하지 않는 파일 형식: {suffix}")
            return None
    
    def get_supported_formats(self) -> List[str]:
        """지원하는 파일 형식 목록"""
        return list(self.supported_formats.keys())


def get_file_parser_service() -> FileParserService:
    """파일 파서 서비스 인스턴스 가져오기"""
    return FileParserService()
